# ABOUTME: Script to generate LunaNMR V1.0 workflow PowerPoint presentation
# ABOUTME: Creates diagram-heavy slides explaining peak picking, fitting, and visualization

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Color scheme
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_BLUE = RGBColor(0x34, 0x98, 0xDB)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)


def set_shape_fill(shape, color):
    """Set solid fill color for shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_shape_line(shape, color, width=Pt(1)):
    """Set line color and width for shape."""
    shape.line.color.rgb = color
    shape.line.width = width


def add_text_to_shape(shape, text, font_size=12, bold=False, color=WHITE):
    """Add text to a shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


def add_box(slide, left, top, width, height, text, fill_color, text_color=WHITE, font_size=11, bold=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Add a colored box with text."""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    set_shape_fill(shape, fill_color)
    set_shape_line(shape, DARK_BLUE, Pt(1))
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.word_wrap = True

    # Vertical centering
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)

    run = shape.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color

    return shape


def add_arrow(slide, start_x, start_y, end_x, end_y, color=DARK_BLUE):
    """Add an arrow connector."""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start_x), Inches(start_y),
        Inches(end_x), Inches(end_y)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector


def add_title(slide, text, font_size=32):
    """Add title to slide."""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.5), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE


def add_subtitle(slide, text, top=0.8):
    """Add subtitle to slide."""
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.5), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(16)
    run.font.color.rgb = GRAY


def create_presentation():
    """Create the LunaNMR workflow presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Get blank layout
    blank_layout = prs.slide_layouts[6]

    # ========== SLIDE 1: Title ==========
    slide1 = prs.slides.add_slide(blank_layout)

    # Main title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "LunaNMR V1.0"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE

    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(1), Inches(3.7), Inches(11.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Advanced NMR Peak Analysis & Integration Workflow"
    run.font.size = Pt(28)
    run.font.color.rgb = LIGHT_BLUE

    # Version info
    ver_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(0.5))
    tf = ver_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "2D Voigt Profile Fitting • PS2D Multi-Peak Deconvolution • Series Integration"
    run.font.size = Pt(14)
    run.font.color.rgb = GRAY

    # ========== SLIDE 2: High-Level Workflow ==========
    slide2 = prs.slides.add_slide(blank_layout)
    add_title(slide2, "Complete Data Integration Workflow")

    # Main flow - vertical cascade
    boxes = [
        ("2D NMR Spectrum\n(Bruker/NMRPipe/SPARKY)", LIGHT_BLUE, 1.3),
        ("Nucleus Detection\n(15N-HSQC / 13C-HSQC)", PURPLE, 2.1),
        ("Peak Detection\n(EnhancedPeakPicker)", GREEN, 2.9),
        ("Overlap Clustering\n(Hierarchical Graph)", ORANGE, 3.7),
        ("Voigt Fitting\n(1D or PS2D 2D)", RED, 4.5),
        ("Quality Assessment\n(R² Categorization)", DARK_BLUE, 5.3),
        ("Output & Visualization", GRAY, 6.1),
    ]

    x_center = 3.0
    for text, color, top in boxes:
        add_box(slide2, x_center, top, 2.8, 0.65, text, color, font_size=10, bold=True)

    # Arrows between boxes
    for i in range(len(boxes) - 1):
        add_arrow(slide2, x_center + 1.4, boxes[i][2] + 0.65, x_center + 1.4, boxes[i+1][2])

    # Side annotations - Detection modes
    add_box(slide2, 6.5, 2.5, 2.2, 0.45, "inplace_mode", GREEN, font_size=9)
    add_box(slide2, 6.5, 3.0, 2.2, 0.45, "full_mode", GREEN, font_size=9)
    add_box(slide2, 6.5, 3.5, 2.2, 0.45, "sn_native", GREEN, font_size=9)

    # Side annotations - Fitting paths
    add_box(slide2, 6.5, 4.3, 2.2, 0.45, "Isolated → 1D", ORANGE, font_size=9)
    add_box(slide2, 6.5, 4.8, 2.2, 0.45, "Overlap → PS2D", ORANGE, font_size=9)

    # Series mode box
    add_box(slide2, 9.2, 2.0, 3.5, 2.0, "Series Modes\n\n• Reference\n• Cascade\n• Detected", PURPLE, font_size=11)

    # Parallel processing box
    add_box(slide2, 9.2, 4.3, 3.5, 1.3, "Processing\n\n• Sequential\n• Parallel (multi-core)", DARK_BLUE, font_size=11)

    # ========== SLIDE 3: Peak Detection ==========
    slide3 = prs.slides.add_slide(blank_layout)
    add_title(slide3, "Peak Detection: EnhancedPeakPicker")

    # Flow diagram
    add_box(slide3, 0.5, 1.5, 2.5, 0.7, "Raw 2D Spectrum", LIGHT_BLUE, font_size=11, bold=True)
    add_arrow(slide3, 3.0, 1.85, 3.5, 1.85)

    add_box(slide3, 3.5, 1.5, 2.5, 0.7, "Noise Estimation\n(Corner IQR)", GREEN, font_size=10)
    add_arrow(slide3, 6.0, 1.85, 6.5, 1.85)

    add_box(slide3, 6.5, 1.5, 2.5, 0.7, "Adaptive\nSmoothing", ORANGE, font_size=10)
    add_arrow(slide3, 9.0, 1.85, 9.5, 1.85)

    add_box(slide3, 9.5, 1.5, 3.0, 0.7, "Local Maxima\nDetection", RED, font_size=10)

    # Second row
    add_box(slide3, 0.5, 2.8, 2.5, 0.7, "Prominence\nAnalysis", PURPLE, font_size=10)
    add_arrow(slide3, 3.0, 3.15, 3.5, 3.15)

    add_box(slide3, 3.5, 2.8, 2.5, 0.7, "Centroid\nRefinement", GREEN, font_size=10)
    add_arrow(slide3, 6.0, 3.15, 6.5, 3.15)

    add_box(slide3, 6.5, 2.8, 2.5, 0.7, "Peak List\nGeneration", DARK_BLUE, font_size=10, bold=True)

    # Detection modes comparison
    add_subtitle(slide3, "Three Detection Modes", top=4.0)

    # Mode boxes
    add_box(slide3, 0.5, 4.5, 3.8, 1.8,
            "inplace_mode\n\n• Reference-based\n• ±X ppm window search\n• Lowest cost\n• Series with known peaks",
            GREEN, font_size=9)

    add_box(slide3, 4.6, 4.5, 3.8, 1.8,
            "full_mode\n\n• De novo detection\n• Complete spectral scan\n• Highest cost\n• Exploratory analysis",
            ORANGE, font_size=9)

    add_box(slide3, 8.7, 4.5, 3.8, 1.8,
            "sn_native\n\n• S/N threshold-based\n• Contour filtering\n• Medium cost\n• Quality filtering",
            PURPLE, font_size=9)

    # ========== SLIDE 4: Overlap Detection & Clustering ==========
    slide4 = prs.slides.add_slide(blank_layout)
    add_title(slide4, "Overlap Detection: Hierarchical Graph Clustering")

    # Elliptical window diagram - left side
    add_subtitle(slide4, "Nucleus-Adaptive Elliptical Windows", top=1.0)

    # Draw ellipses representation as boxes (simplified)
    add_box(slide4, 1.0, 1.5, 2.0, 1.5, "15N\nradF1=0.4 ppm\nradF2=0.04 ppm", LIGHT_BLUE, font_size=9, shape_type=MSO_SHAPE.OVAL)
    add_box(slide4, 3.5, 1.5, 1.5, 1.2, "13C\nradF1=0.15\nradF2=0.04", PURPLE, font_size=8, shape_type=MSO_SHAPE.OVAL)

    # Clustering algorithm
    add_subtitle(slide4, "Clustering Algorithm", top=3.2)

    # Step boxes
    steps = [
        ("1. Build pairwise\noverlap map", GREEN),
        ("2. Transitive\nclosure", ORANGE),
        ("3. Disjoint\npartitions", RED),
        ("4. Route to\nfitter", DARK_BLUE),
    ]

    for i, (text, color) in enumerate(steps):
        add_box(slide4, 0.5 + i*3.2, 3.7, 2.8, 0.8, text, color, font_size=10)
        if i < len(steps) - 1:
            add_arrow(slide4, 3.3 + i*3.2, 4.1, 3.7 + i*3.2, 4.1)

    # Example clustering
    add_subtitle(slide4, "Example: 6 Peaks → 3 Clusters", top=5.0)

    # Peak boxes
    add_box(slide4, 0.5, 5.5, 0.8, 0.5, "A", LIGHT_BLUE, font_size=10, bold=True)
    add_box(slide4, 1.5, 5.5, 0.8, 0.5, "B", LIGHT_BLUE, font_size=10, bold=True)
    add_box(slide4, 2.5, 5.5, 0.8, 0.5, "C", LIGHT_BLUE, font_size=10, bold=True)
    add_box(slide4, 4.0, 5.5, 0.8, 0.5, "D", GREEN, font_size=10, bold=True)
    add_box(slide4, 5.0, 5.5, 0.8, 0.5, "E", GREEN, font_size=10, bold=True)
    add_box(slide4, 6.5, 5.5, 0.8, 0.5, "F", ORANGE, font_size=10, bold=True)

    # Cluster labels
    add_box(slide4, 0.5, 6.2, 2.8, 0.5, "Cluster 1: {A,B,C} → 2D", LIGHT_BLUE, font_size=9)
    add_box(slide4, 4.0, 6.2, 1.8, 0.5, "Cluster 2: {D,E} → 2D", GREEN, font_size=9)
    add_box(slide4, 6.5, 6.2, 1.8, 0.5, "Cluster 3: {F} → 1D", ORANGE, font_size=9)

    # Key point
    key_box = slide4.shapes.add_textbox(Inches(9), Inches(4.5), Inches(3.5), Inches(2.5))
    tf = key_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "KEY: Each peak in exactly ONE cluster\n\nTransitive: A↔B, B↔C\n⟹ {A,B,C} together"
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_BLUE

    # ========== SLIDE 5: Voigt Fitting Decision Tree ==========
    slide5 = prs.slides.add_slide(blank_layout)
    add_title(slide5, "Voigt Fitting: Two-Level Strategy")

    # Decision tree
    add_box(slide5, 5.0, 1.3, 3.3, 0.7, "Peak Cluster", DARK_BLUE, font_size=12, bold=True)

    # Decision diamond (using rounded rect)
    add_box(slide5, 5.0, 2.3, 3.3, 0.7, "Cluster Size?", ORANGE, font_size=11, shape_type=MSO_SHAPE.DIAMOND)

    add_arrow(slide5, 6.65, 2.0, 6.65, 2.3)

    # Two branches
    # Left: Size = 1
    add_box(slide5, 1.5, 3.5, 2.5, 0.6, "Size = 1", GREEN, font_size=10)
    add_box(slide5, 1.5, 4.3, 2.5, 1.5, "1D Cross-Section\nFitting\n\n• F1 slice extraction\n• scipy curve_fit\n• Fast (~0.1s/peak)", GREEN, font_size=9)

    # Right: Size > 1
    add_box(slide5, 9.0, 3.5, 2.5, 0.6, "Size > 1", RED, font_size=10)
    add_box(slide5, 9.0, 4.3, 2.5, 1.5, "PS2D 2D\nSimultaneous\n\n• Multi-peak model\n• 5-stage LM\n• Accurate overlap", RED, font_size=9)

    # Arrows from diamond
    add_arrow(slide5, 5.0, 2.65, 2.75, 3.5)
    add_arrow(slide5, 8.3, 2.65, 10.25, 3.5)

    # Voigt profile equation box
    add_box(slide5, 3.5, 6.0, 6.3, 1.2,
            "Voigt Profile: V(x) = Re[w(z)] / (σ√2π)\nz = (x - pos + iγ) / (σ√2)\n8 params/peak: pos_f1, pos_f2, lw_lor_f1, lw_gau_f1, lw_lor_f2, lw_gau_f2, A, baseline",
            GRAY, font_size=9, text_color=WHITE)

    # ========== SLIDE 6: PS2D 5-Stage Optimization ==========
    slide6 = prs.slides.add_slide(blank_layout)
    add_title(slide6, "PS2D: Multi-Stage Levenberg-Marquardt Optimization")

    # Stage boxes - horizontal flow
    stages = [
        ("Stage 0\nDISABLED\n\n(Was intensity-only\nbut caused collapse)", GRAY),
        ("Stage 1\nLinewidths +\nIntensity\n\nFix: positions\nOpt: lw, A", GREEN),
        ("Stage 2\nPosition\nRefinement\n\nOpt: pos_f1, pos_f2\n(if allowed)", ORANGE),
        ("Stage 4\nGlobal\nRefinement\n\nOpt: ALL params\n(respects flags)", RED),
    ]

    for i, (text, color) in enumerate(stages):
        add_box(slide6, 0.3 + i*2.55, 1.5, 2.4, 2.0, text, color, font_size=9)
        if i < len(stages) - 1:
            add_arrow(slide6, 2.7 + i*2.55, 2.5, 2.85 + i*2.55, 2.5)

    # Convergence diagram
    add_subtitle(slide6, "Convergence Criteria", top=3.8)

    criteria = [
        ("R² > 0.2\n(pragmatic)", GREEN),
        ("χ² reduction\n> 100×", ORANGE),
        ("Formal\nconvergence", RED),
        ("500 iter/stage\n~1500 total", GRAY),
    ]

    for i, (text, color) in enumerate(criteria):
        add_box(slide6, 0.5 + i*3.2, 4.3, 2.8, 0.9, text, color, font_size=10)

    # Advanced features
    add_subtitle(slide6, "Advanced Features", top=5.5)

    features = [
        ("L/G ratio\nconstraint\n(tooclose peaks)", LIGHT_BLUE),
        ("Intensity ratio\npenalty\n(soft constraint)", GREEN),
        ("Adaptive bounds\nfrom PASS1 stats", ORANGE),
        ("Numba JIT\n3-5× speedup", PURPLE),
    ]

    for i, (text, color) in enumerate(features):
        add_box(slide6, 0.5 + i*3.2, 6.0, 2.8, 0.9, text, color, font_size=10)

    # ========== SLIDE 7: Series Integration Modes ==========
    slide7 = prs.slides.add_slide(blank_layout)
    add_title(slide7, "Series Integration: Three Peak Source Modes")

    # Mode 1: Reference
    add_box(slide7, 0.3, 1.5, 4.0, 3.0,
            "REFERENCE Mode\n(Default)\n\n[Spec 1] Detect → Fit → Learn\n        ↓ positions\n[Spec 2+] Search ±X ppm\n         ↓\n      Retain if no peak\n\nUse: Time series, Relaxation",
            LIGHT_BLUE, font_size=10)

    # Mode 2: Detected
    add_box(slide7, 4.6, 1.5, 4.0, 3.0,
            "DETECTED Mode\n\n\n[Each Spectrum]\n    ↓\nDe novo detection\n    ↓\nIndependent fitting\n\n\nUse: Exploratory comparison",
            GREEN, font_size=10)

    # Mode 3: Cascade
    add_box(slide7, 9.0, 1.5, 4.0, 3.0,
            "CASCADE Mode\n\n\n[Spec 1] → positions₁\n    ↓\n[Spec 2] detect near pos₁\n    ↓ positions₂\n[Spec 3] detect near pos₂\n\nUse: Smooth dynamics",
            ORANGE, font_size=10)

    # Series options
    add_subtitle(slide7, "Series Processing Options", top=4.8)

    options = [
        ("lock_cluster_\nassignments\n\nSame clusters\nfor all spectra", PURPLE),
        ("use_ps2d_\nlinewidth_reuse\n\nSpec 1 linewidths\nfor all", DARK_BLUE),
        ("rerun_adaptive_\nper_spectrum\n\nLearn stats\neach spectrum", RED),
    ]

    for i, (text, color) in enumerate(options):
        add_box(slide7, 0.5 + i*4.3, 5.3, 3.8, 1.8, text, color, font_size=9)

    # ========== SLIDE 8: Two-Pass Fitting Strategy ==========
    slide8 = prs.slides.add_slide(blank_layout)
    add_title(slide8, "Two-Pass Fitting Strategy")

    # PASS 1
    add_box(slide8, 0.5, 1.5, 5.5, 2.8,
            "PASS 1: Isolated Peak Learning\n\n"
            "• Fit ONLY non-overlapping peaks\n"
            "• Learn spectrum statistics:\n"
            "  - lw_f1_median (15N/13C linewidth)\n"
            "  - lw_f2_median (1H linewidth)\n"
            "  - α (Gaussian/Lorentzian ratio)\n"
            "• Purpose: Uncontaminated estimates",
            GREEN, font_size=11)

    # Arrow
    add_arrow(slide8, 6.0, 2.9, 6.8, 2.9)

    # PASS 2
    add_box(slide8, 7.0, 1.5, 5.8, 2.8,
            "PASS 2: Full Cluster Fitting\n\n"
            "• Isolated peaks: Use learned stats\n"
            "• Overlap clusters: PS2D 2D fitting\n"
            "  - Skip 1D cross-section (contaminated)\n"
            "  - Use PASS 1 linewidths as initial\n"
            "• All peaks fitted with consistent params",
            ORANGE, font_size=11)

    # Processing modes
    add_subtitle(slide8, "Parallel vs Sequential Processing", top=4.6)

    add_box(slide8, 0.5, 5.1, 6.0, 2.0,
            "Sequential Mode\n\n"
            "• Process clusters one-by-one\n"
            "• Deterministic results\n"
            "• ~120s for 150 peaks",
            LIGHT_BLUE, font_size=10)

    add_box(slide8, 6.8, 5.1, 6.0, 2.0,
            "Parallel Mode (Multi-Core)\n\n"
            "• Different clusters on different cores\n"
            "• IDENTICAL clustering algorithm\n"
            "• ~45s for 150 peaks (2.7× speedup)",
            PURPLE, font_size=10)

    # ========== SLIDE 9: 3D Visualization ==========
    slide9 = prs.slides.add_slide(blank_layout)
    add_title(slide9, "3D Voigt Analysis Visualization")

    # 2x2 layout diagram
    add_subtitle(slide9, "VoigtAnalysisPlotter: 2×2 Axes Grid", top=1.0)

    # Grid representation
    add_box(slide9, 0.5, 1.5, 3.0, 2.0, "ax_x\n\nF2 (1H)\nCross-Section", LIGHT_BLUE, font_size=10)
    add_box(slide9, 3.7, 1.5, 3.0, 2.0, "ax_y\n\nF1 (15N)\nCross-Section", GREEN, font_size=10)
    add_box(slide9, 0.5, 3.7, 3.0, 2.0, "ax_2d\n\n2D Contour\nOverlay", ORANGE, font_size=10)
    add_box(slide9, 3.7, 3.7, 3.0, 2.0, "ax_residuals\n\nResidual\nMap", RED, font_size=10)

    # Layer toggles
    add_subtitle(slide9, "Layer Toggles", top=1.0)
    layers_box = slide9.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5.5), Inches(2.5))
    tf = layers_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "☑ show_experimental (raw spectrum)\n☑ show_fitted (total model)\n☐ show_residuals\n☐ show_individual_peaks\n☑ show_peak_labels"
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_BLUE

    # Color presets
    add_subtitle(slide9, "Color Presets", top=4.2)

    presets = [
        ("Classic\nBlack bg\nSilver data", GRAY),
        ("Clean\nWhite bg\nBlack data", WHITE, DARK_BLUE),
        ("Dark\n#2C3E50 bg\nSilver data", DARK_BLUE),
        ("Warm\nBrown bg\nOrange data", ORANGE),
    ]

    for i, preset in enumerate(presets):
        text = preset[0]
        fill = preset[1]
        text_col = preset[2] if len(preset) > 2 else WHITE
        add_box(slide9, 7.5 + i*1.4, 4.7, 1.3, 1.2, text, fill, text_color=text_col, font_size=8)

    # 3D surface features
    add_box(slide9, 7.5, 6.1, 5.5, 1.0,
            "3D Surface: Wireframe + Surface combined • Peak-specific coloring • Interactive zoom/pan",
            PURPLE, font_size=10)

    # ========== SLIDE 10: Quality Assessment & Output ==========
    slide10 = prs.slides.add_slide(blank_layout)
    add_title(slide10, "Quality Assessment & Output")

    # R² categories
    add_subtitle(slide10, "Quality Categories (R² based)", top=1.0)

    categories = [
        ("Excellent\nR² ≥ 0.95", GREEN),
        ("Good\nR² ≥ 0.85", LIGHT_BLUE),
        ("Fair\nR² ≥ 0.70", ORANGE),
        ("Poor\nR² < 0.70", RED),
    ]

    for i, (text, color) in enumerate(categories):
        add_box(slide10, 0.5 + i*3.2, 1.5, 2.8, 1.0, text, color, font_size=11, bold=True)

    # Output parameters
    add_subtitle(slide10, "Per-Peak Output Parameters", top=2.8)

    params_box = slide10.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(6.0), Inches(3.0))
    tf = params_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = ("• assignment (residue ID)\n"
                "• peak_position (f1, f2 in ppm)\n"
                "• pos_f1, pos_f2 (fitted centers)\n"
                "• lw_lor_f1, lw_gau_f1 (15N linewidths)\n"
                "• lw_lor_f2, lw_gau_f2 (1H linewidths)\n"
                "• amplitude, intensity, volume\n"
                "• r_squared, chi2, iterations\n"
                "• method ('1d' or '2d_simultaneous')\n"
                "• overlap_group_size")
    run.font.size = Pt(11)
    run.font.color.rgb = DARK_BLUE

    # Series output
    add_subtitle(slide10, "Series Analysis Output", top=2.8)

    series_box = slide10.shapes.add_textbox(Inches(7.0), Inches(3.3), Inches(5.5), Inches(3.0))
    tf = series_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = ("• Intensity time series per peak\n"
                "• Error propagation\n"
                "• Decay curve fitting (T1/T1ρ)\n"
                "• Detection rates across series\n"
                "• Quality statistics\n"
                "• CSV export with all parameters")
    run.font.size = Pt(11)
    run.font.color.rgb = DARK_BLUE

    # Performance metrics
    add_subtitle(slide10, "Performance Benchmarks (150 peaks, 15N-HSQC)", top=6.0)

    perf = [
        ("Sequential\n~120s", GRAY),
        ("Parallel (6 cores)\n~45s", PURPLE),
        ("With Numba\n3-5× faster", GREEN),
    ]

    for i, (text, color) in enumerate(perf):
        add_box(slide10, 0.5 + i*4.3, 6.5, 3.8, 0.8, text, color, font_size=10)

    # ========== SLIDE 11: Complete Automated Pipeline ==========
    slide11 = prs.slides.add_slide(blank_layout)
    add_title(slide11, "Complete Automated Pipeline: Data → Series Results")

    # Main pipeline flow - left side vertical
    pipeline_steps = [
        ("1. LOAD\n\nBruker/NMRPipe\nSPARKY formats", LIGHT_BLUE, 1.3),
        ("2. DETECT\n\nAuto nucleus type\n(15N/13C)", GREEN, 2.4),
        ("3. PICK\n\nPeak detection\n+ centroid refine", ORANGE, 3.5),
        ("4. CLUSTER\n\nHierarchical\noverlap groups", PURPLE, 4.6),
        ("5. FIT\n\nPASS1→PASS2\n1D or PS2D 2D", RED, 5.7),
    ]

    for text, color, top in pipeline_steps:
        add_box(slide11, 0.3, top, 2.0, 1.0, text, color, font_size=9, bold=True)

    # Arrows between pipeline steps
    for i in range(len(pipeline_steps) - 1):
        add_arrow(slide11, 1.3, pipeline_steps[i][2] + 1.0, 1.3, pipeline_steps[i+1][2])

    # Series loop
    add_box(slide11, 2.8, 3.0, 2.3, 3.5,
            "SERIES LOOP\n\n┌─────────┐\n│ Spec 1  │→ Learn\n│ Spec 2  │→ Reuse\n│ Spec 3  │→ Reuse\n│   ...   │\n│ Spec N  │→ Reuse\n└─────────┘\n\nPropagate:\n• Positions\n• Linewidths\n• Clusters",
            DARK_BLUE, font_size=8)

    # Unique automation features - right side
    add_subtitle(slide11, "Unique Automation Features", top=1.0)

    features = [
        ("AUTO NUCLEUS\nDETECTION\n\n15N vs 13C from\nspectral range\n→ adaptive params", LIGHT_BLUE),
        ("TWO-PASS\nSTRATEGY\n\nPASS1: Learn from\nisolated peaks\nPASS2: Apply to all", GREEN),
        ("ADAPTIVE\nBOUNDS\n\nmedian + α×MAD\nfrom PASS1 stats\n→ data-driven", ORANGE),
    ]

    for i, (text, color) in enumerate(features):
        add_box(slide11, 5.5 + i*2.6, 1.5, 2.4, 2.0, text, color, font_size=8)

    features2 = [
        ("HIERARCHICAL\nCLUSTERING\n\nTransitive overlap\nDisjoint groups\nOnce per spectrum", PURPLE),
        ("PS2D 2D FIT\n\nSimultaneous\nmulti-peak Voigt\nNo manual tuning", RED),
        ("SOFT\nCONSTRAINTS\n\nL/G ratio penalty\nIntensity guidance\nfor ambiguous peaks", GRAY),
    ]

    for i, (text, color) in enumerate(features2):
        add_box(slide11, 5.5 + i*2.6, 3.8, 2.4, 2.0, text, color, font_size=8)

    # Output arrow and final results
    add_box(slide11, 5.5, 6.2, 7.3, 1.0,
            "OUTPUT: Per-peak intensities × N spectra  •  Decay curves  •  Quality metrics  •  CSV export",
            DARK_BLUE, font_size=10, bold=True)

    # Key differentiator box
    key_box = slide11.shapes.add_textbox(Inches(0.3), Inches(6.8), Inches(4.8), Inches(0.5))
    tf = key_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Zero manual parameter tuning required"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RED

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), "LunaNMR_V1_Workflow.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
